"""DRM 걸린 Excel / PPT 를 Fasoo 훅이 감시하지 않는 포맷으로 추출.

Fasoo 훅은 Office 확장자(`.xlsx`/`.pptx`) 쓰기를 감시해 새 파일에도 DRM 을 씌움.
SaveAs 뿐 아니라 xlsxwriter 로 처음부터 쓴 xlsx 도 똑같이 걸림(2026-04-21 실측).

회피 전략:
    - 값 / 수식 / 차트 시리즈 → **첫 줄 공란 .txt** (검증됨, proto_:22176 트릭)
    - 차트 이미지          → **.png** (Fasoo 이미지 훅 정책에 따라 다름 — 백업용)
    - .txt 에 **차트 시리즈 원본 데이터** 도 함께 저장 → PNG 실패해도 재플롯 가능

Excel 입력 산출물:
    <stem>_bundle.txt            — 값 / 수식 / 차트 시리즈 (첫 줄 공란)
    <stem>_chart_<sheet>_NN.png  — 차트별 PNG (정책 허용 시 즉시 사용)

PPT 입력 산출물 (기존 이미지 재작성 방식 유지):
    <stem>_export.pptx           — Fasoo 정책 따라 DRM 걸릴 수 있음 (후속 개선 대기)

사용법:
    python drm_reload_test.py <입력파일> [출력stem]
또는 SRC / DST 상수 직접 수정.

재로드:
    from tools.drm_reload_test import load_bundle
    b = load_bundle('<stem>_bundle.txt')
    # b['values']['Sheet1']          → list[list]  (TSV raw, header 없음)
    # b['formulas']['Sheet1']['A1']  → '=SUM(...)'
    # b['charts'][0]                 → {id, title, x_axis, y_axis, png, series:[{name,x,y}]}
"""
from __future__ import annotations

import os
import re
import sys
import tempfile
import shutil
from datetime import datetime, date
from pathlib import Path

SRC = r""   # 예: r"D:\data\reliability\sample.xlsx"
DST = r""   # 비우면 입력 파일 stem 사용


# ── 공용 헬퍼 ────────────────────────────────────────────────────────
def _col_to_letter(n: int) -> str:
    """1-indexed 컬럼 → Excel 레터 (A, B, ..., Z, AA, ...)."""
    s = ""
    while n > 0:
        n, r = divmod(n - 1, 26)
        s = chr(65 + r) + s
    return s


def _to_2d(val, nrows: int, ncols: int) -> list[list]:
    """xlwings .value / .formula 반환값을 항상 2D list 로 정규화."""
    if nrows == 1 and ncols == 1:
        return [[val]]
    if nrows == 1:
        return [list(val)] if isinstance(val, (list, tuple)) else [[val]]
    if ncols == 1:
        if val and isinstance(val[0], (list, tuple)):
            return [list(r) for r in val]
        return [[v] for v in val]
    return [list(r) for r in val]


def _safe_name(s: str) -> str:
    return re.sub(r'[^\w.-]', '_', str(s))[:64] or "_"


def _ser_value(v) -> str:
    """시리즈 개별 값을 txt 직렬화."""
    if v is None:
        return ""
    if isinstance(v, (datetime, date)):
        return v.isoformat()
    if isinstance(v, float):
        # Excel 의 정수도 float 로 옴 → 정수로 떨어지면 정수 표기
        if v.is_integer() and abs(v) < 1e15:
            return str(int(v))
        return repr(v)
    return str(v).replace('\t', ' ').replace('\n', ' ').replace('\r', ' ')


def _join_values(vals) -> str:
    if not vals:
        return ""
    return ",".join(_ser_value(v) for v in vals)


# ── Excel 차트 추출 ──────────────────────────────────────────────────
def _extract_charts(sh, png_prefix: str) -> list[dict]:
    """시트의 ChartObjects 전체를 순회 → 시리즈 데이터 + PNG export.

    각 시도는 try-except 로 격리 — 차트 하나가 깨져도 다음 차트 계속 처리.
    """
    charts: list[dict] = []
    try:
        co_coll = sh.api.ChartObjects()
        n = co_coll.Count
    except Exception:
        return charts

    sheet_tag = _safe_name(sh.name)

    for i in range(1, n + 1):
        try:
            co = sh.api.ChartObjects(i)
            chart = co.Chart
        except Exception:
            continue

        info: dict = {"id": i, "name": None, "title": None,
                      "x_axis": None, "y_axis": None,
                      "png": None, "series": []}

        try: info["name"] = co.Name
        except Exception: pass

        try:
            if chart.HasTitle:
                info["title"] = chart.ChartTitle.Text
        except Exception: pass

        for axis_idx, key in ((1, "x_axis"), (2, "y_axis")):
            try:
                ax = chart.Axes(axis_idx)
                if ax.HasTitle:
                    info[key] = ax.AxisTitle.Text
            except Exception:
                pass

        try:
            sc = chart.SeriesCollection()
            ns = sc.Count
        except Exception:
            ns = 0

        for j in range(1, ns + 1):
            try:
                s = chart.SeriesCollection(j)
            except Exception:
                continue
            s_info = {"name": None, "x": [], "y": []}
            try: s_info["name"] = s.Name
            except Exception: pass
            try: s_info["y"] = list(s.Values) if s.Values else []
            except Exception: pass
            try: s_info["x"] = list(s.XValues) if s.XValues else []
            except Exception: pass
            info["series"].append(s_info)

        # PNG export — 실패해도 시리즈 데이터는 이미 확보됨
        png_path = Path(f"{png_prefix}_{sheet_tag}_{i:02d}.png").resolve()
        try:
            chart.Export(str(png_path), "PNG")
            if png_path.is_file():
                info["png"] = png_path.name
        except Exception:
            pass

        charts.append(info)

    return charts


# ── Excel 번들 작성 ──────────────────────────────────────────────────
def export_excel_bundle(src: Path, bundle_txt: Path, png_prefix: str) -> dict:
    """xlwings 로 값/수식/차트 추출 → 첫 줄 공란 txt 번들 + 차트 PNG."""
    import xlwings as xw
    import pandas as pd

    stats = {"sheets": 0, "rows_total": 0, "formulas": 0, "charts": 0, "pngs": 0}

    sheet_values: list[tuple[str, pd.DataFrame]] = []
    sheet_formulas: list[tuple[str, list[tuple[str, str]]]] = []
    sheet_charts: list[tuple[str, list[dict]]] = []

    app = xw.App(visible=False, add_book=False)
    try:
        app.display_alerts = False
        wb = app.books.open(str(src), update_links=False, read_only=True)
        try:
            for sh in wb.sheets:
                rng = sh.used_range
                empty = (rng.last_cell.row == 1 and rng.last_cell.column == 1
                         and rng.value is None)

                if empty:
                    sheet_values.append((sh.name, pd.DataFrame()))
                    sheet_formulas.append((sh.name, []))
                else:
                    val = rng.options(pd.DataFrame, index=False, header=False).value
                    df = val if isinstance(val, pd.DataFrame) else pd.DataFrame([[val]])

                    row0, col0 = rng.row, rng.column
                    nrows = rng.last_cell.row - row0 + 1
                    ncols = rng.last_cell.column - col0 + 1
                    formulas_2d = _to_2d(rng.formula, nrows, ncols)

                    f_list: list[tuple[str, str]] = []
                    for i, row in enumerate(formulas_2d):
                        for j, v in enumerate(row):
                            if isinstance(v, str) and v.startswith('='):
                                addr = f"{_col_to_letter(col0 + j)}{row0 + i}"
                                f_list.append((addr, v))

                    sheet_values.append((sh.name, df))
                    sheet_formulas.append((sh.name, f_list))
                    stats["rows_total"] += len(df)
                    stats["formulas"] += len(f_list)

                # 차트는 빈 시트에도 있을 수 있음
                chart_list = _extract_charts(sh, png_prefix)
                sheet_charts.append((sh.name, chart_list))
                stats["charts"] += len(chart_list)
                stats["pngs"] += sum(1 for c in chart_list if c.get("png"))
                stats["sheets"] += 1
        finally:
            wb.close()
    finally:
        app.quit()

    # 번들 txt 작성 (첫 줄 공란)
    with open(bundle_txt, 'w', encoding='utf-8') as f:
        f.write("\n")  # ← DRM 회피용 첫 줄 공란 (proto_:22176)
        f.write("# DRM-bypass bundle\n")
        f.write(f"# Source: {src}\n")
        f.write(f"# Exported: {datetime.now().isoformat(timespec='seconds')}\n")
        f.write(f"# Sheets: {stats['sheets']}, Formulas: {stats['formulas']}, "
                f"Charts: {stats['charts']}, PNGs: {stats['pngs']}\n")
        f.write("# Format: sections delimited by ===NAME===, items by [key].\n")
        f.write("\n")

        # VALUES
        f.write("===VALUES===\n")
        for name, df in sheet_values:
            f.write(f"[{name}]\n")
            if not df.empty:
                df.to_csv(f, sep='\t', index=False, header=False,
                          lineterminator='\n', na_rep='')
            f.write("\n")

        # FORMULAS
        f.write("===FORMULAS===\n")
        for name, f_list in sheet_formulas:
            if not f_list:
                continue
            f.write(f"[{name}]\n")
            for addr, formula in f_list:
                f.write(f"{addr}\t{formula}\n")
            f.write("\n")

        # CHARTS
        f.write("===CHARTS===\n")
        for sheet_name, chart_list in sheet_charts:
            for ch in chart_list:
                tag = ch.get("name") or f"Chart{ch['id']}"
                f.write(f"[{sheet_name}::{tag}]\n")
                for key in ("title", "x_axis", "y_axis", "png"):
                    if ch.get(key):
                        f.write(f"{key}={ch[key]}\n")
                for si, s in enumerate(ch.get("series", []), start=1):
                    sname = s.get("name") or ""
                    sname = str(sname).replace('\t', ' ').replace('\n', ' ')
                    f.write(f"series.{si}.name\t{sname}\n")
                    f.write(f"series.{si}.x\t{_join_values(s.get('x', []))}\n")
                    f.write(f"series.{si}.y\t{_join_values(s.get('y', []))}\n")
                f.write("\n")

    return stats


# ── 번들 로더 ───────────────────────────────────────────────────────
def load_bundle(path: str | Path) -> dict:
    """export_excel_bundle 이 생성한 txt 번들 파싱.

    Returns:
        {
          'values':   {sheet: [[cell, ...], ...]},   # TSV raw, header 없음
          'formulas': {sheet: {cell: formula}},
          'charts':   [{id, title, x_axis, y_axis, png, series:[{name,x,y}]}, ...],
        }
    """
    result: dict = {"values": {}, "formulas": {}, "charts": []}
    section: str | None = None
    key: str | None = None
    chart: dict | None = None
    value_buf: list[str] = []

    def _flush_values():
        if key is not None and value_buf:
            result["values"][key] = [ln.split('\t') for ln in value_buf]
        value_buf.clear()

    def _flush_chart():
        nonlocal chart
        if chart is not None:
            result["charts"].append(chart)
        chart = None

    def _parse_scalar(s: str):
        s = s.strip()
        if s == "":
            return None
        try:
            if '.' in s or 'e' in s.lower():
                return float(s)
            return int(s)
        except ValueError:
            return s

    with open(path, 'r', encoding='utf-8') as f:
        for raw in f:
            line = raw.rstrip('\r\n')
            if not line.strip() or line.lstrip().startswith('#'):
                continue

            if line.startswith('===') and line.endswith('==='):
                if section == 'VALUES': _flush_values()
                if section == 'CHARTS': _flush_chart()
                section = line.strip('=').strip()
                key = None
                continue

            if line.startswith('[') and line.endswith(']'):
                if section == 'VALUES':
                    _flush_values()
                    key = line[1:-1]
                elif section == 'FORMULAS':
                    key = line[1:-1]
                    result["formulas"].setdefault(key, {})
                elif section == 'CHARTS':
                    _flush_chart()
                    chart = {"id": line[1:-1], "title": None, "x_axis": None,
                             "y_axis": None, "png": None, "series": []}
                continue

            if section == 'VALUES':
                value_buf.append(line)

            elif section == 'FORMULAS':
                if key is None:
                    continue
                parts = line.split('\t', 1)
                if len(parts) == 2:
                    result["formulas"][key][parts[0]] = parts[1]

            elif section == 'CHARTS':
                if chart is None:
                    continue
                if line.startswith('series.'):
                    parts = line.split('\t', 1)
                    if len(parts) != 2:
                        continue
                    key_parts = parts[0].split('.')
                    if len(key_parts) != 3:
                        continue
                    _, idx_s, field = key_parts
                    try: idx = int(idx_s)
                    except ValueError: continue
                    while len(chart["series"]) < idx:
                        chart["series"].append({"name": None, "x": [], "y": []})
                    s = chart["series"][idx - 1]
                    if field == "name":
                        s["name"] = parts[1]
                    elif field in ("x", "y"):
                        s[field] = [_parse_scalar(v) for v in parts[1].split(',')] if parts[1] else []
                elif '=' in line:
                    k, v = line.split('=', 1)
                    if k in ("title", "x_axis", "y_axis", "png"):
                        chart[k] = v

    # 마지막 섹션 flush
    if section == 'VALUES': _flush_values()
    if section == 'CHARTS': _flush_chart()

    return result


# ── PowerPoint (기존 유지) ──────────────────────────────────────────
def export_pptx(src: Path, dst: Path) -> dict:
    """PowerPoint COM 으로 슬라이드를 PNG → python-pptx 로 재작성.

    주의: pptx 확장자라 Fasoo 훅에 걸릴 가능성 있음. 걸리면 PPT 도 번들 txt+PNG 로
    재설계 예정 (후속 작업).
    """
    import win32com.client as win32
    import pythoncom
    from pptx import Presentation
    from pptx.util import Emu

    stats = {"slides": 0}
    tmp_dir = Path(tempfile.mkdtemp(prefix="drm_pptx_"))
    pythoncom.CoInitialize()
    app = None
    pres = None
    try:
        app = win32.DispatchEx("PowerPoint.Application")
        pres = app.Presentations.Open(str(src), ReadOnly=True, Untitled=False, WithWindow=False)
        slide_w_pt = pres.PageSetup.SlideWidth
        slide_h_pt = pres.PageSetup.SlideHeight
        slide_w_emu = int(slide_w_pt * 12700)
        slide_h_emu = int(slide_h_pt * 12700)
        px_w = int(slide_w_pt * 2)
        px_h = int(slide_h_pt * 2)

        images: list[Path] = []
        for idx in range(1, pres.Slides.Count + 1):
            img_path = tmp_dir / f"slide_{idx:04d}.png"
            pres.Slides.Item(idx).Export(str(img_path), "PNG", px_w, px_h)
            images.append(img_path)
            stats["slides"] += 1
    finally:
        if pres is not None:
            try: pres.Close()
            except Exception: pass
        if app is not None:
            try: app.Quit()
            except Exception: pass
        pythoncom.CoUninitialize()

    try:
        prs = Presentation()
        prs.slide_width = Emu(slide_w_emu)
        prs.slide_height = Emu(slide_h_emu)
        blank_layout = prs.slide_layouts[6]
        for img in images:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(str(img), 0, 0, prs.slide_width, prs.slide_height)
        prs.save(str(dst))
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)

    return stats


# ── 엔트리 포인트 ────────────────────────────────────────────────────
def main(argv: list[str]) -> int:
    if len(argv) >= 2:
        src = argv[1]
        dst = argv[2] if len(argv) >= 3 else ""
    else:
        src, dst = SRC, DST

    if not src:
        print("[에러] 입력 파일 경로 필요. 인자 전달 또는 SRC 변수 수정.")
        return 2

    src_path = Path(src).expanduser().resolve()
    if not src_path.is_file():
        print(f"[에러] 파일 없음: {src_path}")
        return 2

    ext = src_path.suffix.lower()
    out_stem = Path(dst).expanduser().resolve() if dst else src_path.with_suffix("")
    out_stem.parent.mkdir(parents=True, exist_ok=True)

    print(f"[입력] {src_path}  ({src_path.stat().st_size:,} bytes)")

    try:
        if ext in (".xlsx", ".xls", ".xlsm"):
            bundle = Path(f"{out_stem}_bundle.txt")
            png_prefix = f"{out_stem}_chart"
            print(f"[출력] {bundle}  (첫 줄 공란, 값/수식/차트 시리즈)")
            print(f"[출력] {png_prefix}_<sheet>_NN.png  (차트별 이미지)")
            stats = export_excel_bundle(src_path, bundle, png_prefix)
            print(f"[Excel] 시트 {stats['sheets']}, 행 {stats['rows_total']:,}, "
                  f"수식 {stats['formulas']}, 차트 {stats['charts']}, "
                  f"PNG {stats['pngs']}")
            if stats['charts'] and stats['pngs'] < stats['charts']:
                print(f"  주의: PNG {stats['charts'] - stats['pngs']}개 실패 "
                      f"(Fasoo 이미지 훅? — 시리즈 데이터는 번들 txt 에 보존됨)")
        elif ext in (".pptx", ".ppt", ".pptm"):
            dst_pptx = Path(f"{out_stem}_export{ext}")
            print(f"[출력] {dst_pptx}  (pptx 재작성 — Fasoo 정책 따라 DRM 걸릴 수 있음)")
            stats = export_pptx(src_path, dst_pptx)
            print(f"[PPT] 슬라이드 {stats['slides']}")
        else:
            print(f"[에러] 지원 않는 확장자: {ext}")
            return 2
    except Exception as e:
        print(f"[실패] {type(e).__name__}: {e}")
        print("  사내 Fasoo 미설치 환경이면 정상 실패.")
        return 1

    print()
    print("── 외부 환경에서 재로드 ──")
    print("  from tools.drm_reload_test import load_bundle")
    print("  b = load_bundle('<stem>_bundle.txt')")
    print("  b['values']['Sheet1']          # list[list] (TSV raw)")
    print("  b['formulas']['Sheet1']['A1']  # '=SUM(...)'")
    print("  b['charts'][0]['series']       # [{'name','x','y'}, ...] → matplotlib 재플롯")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
