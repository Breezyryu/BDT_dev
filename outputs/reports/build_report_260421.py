"""선행Lab 성능/수명 파트 — BDT 현황 및 26년 계획 (그룹장 보고자료) v2

v2 (2026-04-21): 배터리 도메인 지식(Science Note) 보강
  - 미션별 슬라이드에 GITT/dV/dQ/DFN/Arrhenius/Fade model 수식 근거 추가
  - Executive Summary 하단에 과학적 흐름 띠 추가
출력: outputs/reports/260421_보고_그룹장_선행Lab_BDT_현황및계획.pptx
"""

from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─────────────────────────── 색상 팔레트 ────────────────────────────────────
NAVY        = RGBColor(0x1E, 0x27, 0x61)
BLUE        = RGBColor(0x2E, 0x86, 0xAB)
LIGHT_BLUE  = RGBColor(0xCA, 0xDC, 0xFC)
TEAL        = RGBColor(0x02, 0x80, 0x90)
SAND        = RGBColor(0xF5, 0xEE, 0xDC)
RED         = RGBColor(0xC8, 0x10, 0x2E)
GREY_TXT    = RGBColor(0x44, 0x44, 0x44)
GREY_MUTED  = RGBColor(0x88, 0x88, 0x88)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x10, 0x10, 0x10)
GREEN_TAG   = RGBColor(0x2C, 0x5F, 0x2D)
ORANGE_TAG  = RGBColor(0xD9, 0x7C, 0x00)
GREY_TAG    = RGBColor(0x77, 0x77, 0x77)
SCI_YELLOW  = RGBColor(0xFF, 0xD7, 0x00)  # Science Note 강조

FONT_TITLE = "맑은 고딕"
FONT_BODY  = "맑은 고딕"
FONT_MONO  = "Consolas"   # 수식용

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    shp.text_frame.text = ""
    return shp


def add_round_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.12
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=GREY_TXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_BODY,
             italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_confidential_stamp(slide):
    add_rect(slide, Inches(11.65), Inches(0.22), Inches(1.30), Inches(0.28),
             WHITE, line=RED)
    add_text(slide, Inches(11.65), Inches(0.22), Inches(1.30), Inches(0.28),
             "Confidential", size=10, bold=True, color=RED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(10.60), Inches(0.55), Inches(2.35), Inches(0.28),
             "Advanced Battery Lab.", size=10, bold=True, color=BLUE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_title(slide, num, text):
    add_text(slide, Inches(0.45), Inches(0.22), Inches(10.5), Inches(0.75),
             f"{num}. {text}", size=30, bold=True, color=NAVY,
             font=FONT_TITLE, anchor=MSO_ANCHOR.MIDDLE)


def add_footer(slide, page_no, total):
    add_text(slide, Inches(0.45), Inches(7.1), Inches(5.0), Inches(0.25),
             "선행Lab 성능/수명 파트 | BDT 현황 및 26년 계획",
             size=9, color=GREY_MUTED)
    add_text(slide, Inches(11.8), Inches(7.1), Inches(1.2), Inches(0.25),
             f"{page_no} / {total}", size=9, color=GREY_MUTED,
             align=PP_ALIGN.RIGHT)


def add_status_pill(slide, x, y, w, h, label, color):
    add_round_rect(slide, x, y, w, h, fill=color)
    add_text(slide, x, y, w, h, label, size=9, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ─────────────────────────── 프레젠테이션 ────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]
TOTAL = 10


# ============================ S1. 표지 ============================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, Inches(8.5), SLIDE_H, NAVY)
add_rect(s, Inches(8.5), 0, Inches(4.833), SLIDE_H, LIGHT_BLUE)
add_rect(s, 0, Inches(2.0), Inches(0.45), Inches(3.5), RED)

add_rect(s, Inches(11.65), Inches(0.22), Inches(1.30), Inches(0.28),
         WHITE, line=RED)
add_text(s, Inches(11.65), Inches(0.22), Inches(1.30), Inches(0.28),
         "Confidential", size=10, bold=True, color=RED,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.8), Inches(0.7), Inches(7), Inches(0.4),
         "Advanced Battery Lab. | 선행Lab. 성능/수명 파트",
         size=13, color=LIGHT_BLUE, font=FONT_TITLE)

add_text(s, Inches(0.8), Inches(2.3), Inches(7.5), Inches(1.3),
         "BDT 개발 현황 및",
         size=44, bold=True, color=WHITE, font=FONT_TITLE)
add_text(s, Inches(0.8), Inches(3.2), Inches(7.5), Inches(1.3),
         "26년 연간 계획",
         size=44, bold=True, color=WHITE, font=FONT_TITLE)

add_text(s, Inches(0.8), Inches(4.7), Inches(7.5), Inches(0.5),
         "BatteryDataTool — 4대 미션 통합 플랫폼",
         size=20, italic=True, color=LIGHT_BLUE, font=FONT_TITLE)

add_rect(s, Inches(0.8), Inches(6.3), Inches(7.0), Inches(0.03), RED)
add_text(s, Inches(0.8), Inches(6.4), Inches(7.5), Inches(0.35),
         "2026-04-21 | 그룹장 보고",
         size=14, color=WHITE, font=FONT_TITLE)
add_text(s, Inches(0.8), Inches(6.75), Inches(7.5), Inches(0.35),
         "작성 : 류 (성능/수명 파트)",
         size=12, color=LIGHT_BLUE, font=FONT_TITLE)

keywords = [
    ("소재 DB",   "전기화학 물성·OCP·GITT"),
    ("dV/dQ",     "LLI / LAM 분리"),
    ("성능 Sim",  "P2D · SPMe · SPM"),
    ("수명 예측", "Empirical + Arrhenius"),
]
ky = Inches(1.2)
for label, sub in keywords:
    add_round_rect(s, Inches(9.0), ky, Inches(3.8), Inches(1.1), WHITE)
    add_text(s, Inches(9.0), ky + Inches(0.12), Inches(3.8), Inches(0.45),
             label, size=20, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
    add_text(s, Inches(9.0), ky + Inches(0.58), Inches(3.8), Inches(0.45),
             sub, size=11, color=BLUE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
    ky += Inches(1.35)


# ============================ S2. Executive Summary ============================
s = prs.slides.add_slide(BLANK)
add_confidential_stamp(s)
add_title(s, "Executive", "Summary — 한 장 요약")

items = [
    ("선행Lab 4대 미션",       "진행 중",   "소재 DB · dV/dQ · 성능 Sim · 수명 예측", ORANGE_TAG),
    ("BDT 통합 플랫폼",         "1Q 완료",  "4대 미션의 공용 실행 도구로 자리매김",   GREEN_TAG),
    ("4월 개발 실적",           "162 커밋", "파이프라인 리팩터 · 분류 v4 · OSS 14종", TEAL),
    ("2분기 Focus",            "계획 확정", "DB UI · 풀셀↔단극 · 캘리브 · Si 계수",   BLUE),
    ("선결 과제",              "관리 중",   "DRM 포맷 · OSS 라이선스 경계",          RED),
]

ty = Inches(1.3)
for title, tag, desc, tag_color in items:
    add_rect(s, Inches(0.45), ty, Inches(8.5), Inches(0.85), SAND)
    add_rect(s, Inches(0.45), ty, Inches(0.12), Inches(0.85), NAVY)
    add_text(s, Inches(0.75), ty + Inches(0.08), Inches(3.0), Inches(0.40),
             title, size=16, bold=True, color=NAVY, font=FONT_TITLE)
    add_status_pill(s, Inches(0.75), ty + Inches(0.48), Inches(1.3),
                    Inches(0.28), tag, tag_color)
    add_text(s, Inches(2.3), ty + Inches(0.48), Inches(6.3), Inches(0.32),
             desc, size=12, color=GREY_TXT, font=FONT_BODY,
             anchor=MSO_ANCHOR.MIDDLE)
    ty += Inches(0.97)

# 하단 : 과학적 흐름 띠 (도메인 지식)
flow_y = Inches(6.17)
add_rect(s, Inches(0.45), flow_y, Inches(8.5), Inches(0.82), NAVY)
add_text(s, Inches(0.65), flow_y + Inches(0.05), Inches(8.2), Inches(0.30),
         "과학적 흐름 (Science Pipeline)",
         size=11, bold=True, color=SCI_YELLOW, font=FONT_TITLE)

flow_items = [("소재 물성", "U(x), D_s, i₀"),
              ("열화 모드", "LLI · LAM"),
              ("전기화학 모델", "DFN · SPMe"),
              ("수명 예측", "Q(n) = 1 - an^b - ce^d(n-e)")]
fx = Inches(0.65); fw = Inches(1.90)
for i, (t, sub) in enumerate(flow_items):
    x = fx + (fw + Inches(0.08)) * i
    add_round_rect(s, x, flow_y + Inches(0.38), fw, Inches(0.40), LIGHT_BLUE)
    add_text(s, x, flow_y + Inches(0.38), fw, Inches(0.20),
             t, size=10, bold=True, color=NAVY, font=FONT_TITLE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x, flow_y + Inches(0.58), fw, Inches(0.20),
             sub, size=8, italic=True, color=BLUE, font=FONT_MONO,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 화살표
    if i < 3:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    x + fw + Inches(0.005),
                                    flow_y + Inches(0.48),
                                    Inches(0.07), Inches(0.20))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = SCI_YELLOW
        arrow.line.fill.background(); arrow.shadow.inherit = False

# 우측: 핵심 메시지 박스
add_rect(s, Inches(9.2), Inches(1.3), Inches(3.8), Inches(5.7), NAVY)
add_text(s, Inches(9.4), Inches(1.5), Inches(3.4), Inches(0.45),
         "핵심 메시지", size=14, bold=True, color=LIGHT_BLUE, font=FONT_TITLE)
add_rect(s, Inches(9.4), Inches(1.95), Inches(0.45), Inches(0.04), RED)

msg = ("발표자료 슬라이드 2의\n"
       "4대 기둥은\n\n"
       "BDT 라는\n"
       "하나의 도구 위에서\n"
       "통합 구동\n\n"
       "되고 있으며,\n\n"
       "1분기  \"기반 구축\"\n"
       "        ↓\n"
       "2분기  \"결합 · 정밀화\"\n\n"
       "로 자연스럽게 전환 중.")
add_text(s, Inches(9.4), Inches(2.15), Inches(3.4), Inches(4.8),
         msg, size=12, color=WHITE, font=FONT_TITLE)

add_footer(s, 2, TOTAL)


# ============================ S3. 선행Lab 전체 그림 ============================
s = prs.slides.add_slide(BLANK)
add_confidential_stamp(s)
add_title(s, "1", "선행Lab. 성능/수명 파트 — 전체 그림")

missions = [
    ("①", "소재 전기화학 물성 DB화",  "설계 경쟁력 기반 데이터 구축",
     "BDT : 파형수정 · 셋 결과 · DB 연동(예정)"),
    ("②", "dV/dQ 양음극 분리",        "열화 모드 (LLI / LAM) 규명",
     "BDT : dVdQ 분석 탭"),
    ("③", "성능 Simulation (전기화학)", "온도별 profile · 급속충전 risk",
     "BDT : 전기화학 시뮬레이션 탭 (PyBaMM)"),
    ("④", "수명 예측 (Empirical)",    "승인용 / 실사용 수명 예측",
     "BDT : Eu · 승인 · 실사용 수명예측 탭"),
]

positions = [
    (Inches(0.45), Inches(1.15)),
    (Inches(6.90), Inches(1.15)),
    (Inches(0.45), Inches(2.75)),
    (Inches(6.90), Inches(2.75)),
]
for (num, title, goal, bdt), (x, y) in zip(missions, positions):
    add_rect(s, x, y, Inches(6.0), Inches(1.5), LIGHT_BLUE)
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15),
                                y + Inches(0.15), Inches(0.55), Inches(0.55))
    circle.fill.solid(); circle.fill.fore_color.rgb = NAVY
    circle.line.fill.background(); circle.shadow.inherit = False
    add_text(s, x + Inches(0.15), y + Inches(0.15), Inches(0.55), Inches(0.55),
             num, size=20, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
    add_text(s, x + Inches(0.85), y + Inches(0.15), Inches(5.0), Inches(0.35),
             title, size=15, bold=True, color=NAVY, font=FONT_TITLE)
    add_text(s, x + Inches(0.85), y + Inches(0.55), Inches(5.0), Inches(0.35),
             goal, size=11, color=GREY_TXT, font=FONT_BODY)
    add_text(s, x + Inches(0.20), y + Inches(1.05), Inches(5.7), Inches(0.35),
             bdt, size=10, italic=True, color=BLUE, font=FONT_BODY)

# 하단 : 26년 분기별 타임라인
add_text(s, Inches(0.45), Inches(4.45), Inches(10), Inches(0.4),
         "26년 분기별 타임라인", size=16, bold=True, color=NAVY, font=FONT_TITLE)

qx0 = Inches(0.45); qw = Inches(3.10)
qy = Inches(4.9); qh = Inches(1.95)
for i, qlabel in enumerate(["1Q — 기반 구축", "2Q — 결합 · 정밀화",
                              "3Q — 자동화 · 통합", "4Q — 트렌드 · 운영"]):
    x = qx0 + qw * i + Inches(0.05 * i)
    color = NAVY if i < 2 else BLUE
    add_rect(s, x, qy, qw, Inches(0.4), color)
    add_text(s, x, qy, qw, Inches(0.4), qlabel,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)

rows = [
    ("소재 DB",   ["물성 파라미터 설계\n코인/삼전극 SOP",
                    "DB 데이터 관리 UI 연동",
                    "소재 DB 파이프라인 정립",
                    "지속 업데이트"]),
    ("dV/dQ",     ["dVdQ 스무딩\n음극 OCV 모델",
                    "풀셀↔단극 피크 매칭\nLLI/LAM 추출",
                    "Tool 내부 dVdQ 통합\n분리 정확도 개선",
                    "다수 사이클 트렌드\n자동 분석"]),
    ("성능 Sim",  ["오픈소스 모델(UI)\n충/방전·GITT·HPPC",
                    "기본 모델 캘리브\nArrhenius 온도 의존성",
                    "급속 스텝 가상 시나리오\n리튬 석출 risk",
                    "지속 고도화"]),
    ("수명 예측", ["온도별 파라미터 도출\n(folderble)",
                    "Si 하한전압 계수\n승인/실사용 EOL",
                    "다조건 비교 + 자동화\n리포트 산출",
                    "지속 운영"]),
]

ry = qy + Inches(0.45); row_h = Inches(0.36)
for ri, (miss_label, cells) in enumerate(rows):
    for ci, cell_text in enumerate(cells):
        x = qx0 + qw * ci + Inches(0.05 * ci)
        bg = LIGHT_BLUE if ri % 2 == 0 else SAND
        add_rect(s, x, ry + row_h * ri, qw, row_h, bg)
        add_text(s, x + Inches(0.08), ry + row_h * ri, qw, row_h,
                 cell_text, size=8, color=GREY_TXT, font=FONT_BODY,
                 anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 3, TOTAL)


# ============================ S4. BDT 통합 플랫폼 ============================
s = prs.slides.add_slide(BLANK)
add_confidential_stamp(s)
add_title(s, "2", "BDT — 4대 미션 통합 플랫폼")

add_rect(s, Inches(0.45), Inches(1.15), Inches(5.5), Inches(5.6), LIGHT_BLUE)
add_text(s, Inches(0.65), Inches(1.3), Inches(5.2), Inches(0.5),
         "왜 BDT 가 중심에 오는가",
         size=17, bold=True, color=NAVY, font=FONT_TITLE)

add_text(s, Inches(0.7), Inches(1.95), Inches(5.0), Inches(0.35),
         "과거 (개별 스크립트)", size=12, bold=True, color=RED, font=FONT_TITLE)
past_items = [
    "· 각자 Matlab / Python / Excel 로 분석",
    "· 결과 포맷 불일치, 재현성 낮음",
    "· 스크립트 소재 불명, 인수인계 어려움",
]
yy = Inches(2.3)
for it in past_items:
    add_text(s, Inches(0.9), yy, Inches(5.0), Inches(0.28), it,
             size=11, color=GREY_TXT, font=FONT_BODY)
    yy += Inches(0.28)

add_text(s, Inches(0.7), yy + Inches(0.15), Inches(5.0), Inches(0.35),
         "현재 (BDT 단일 GUI)", size=12, bold=True, color=TEAL, font=FONT_TITLE)
now_items = [
    "· raw 로딩 · 분석 · Sim · 수명예측 통합",
    "· 동일 입력 → 동일 경로 → 재현성 확보",
    "· 10개 모델 비교 : 수십 초 → 수 초 (DB 예정)",
    "· 신입 학습 곡선 단축",
]
yy += Inches(0.5)
for it in now_items:
    add_text(s, Inches(0.9), yy, Inches(5.0), Inches(0.28), it,
             size=11, color=GREY_TXT, font=FONT_BODY)
    yy += Inches(0.28)

add_rect(s, Inches(0.65), Inches(5.65), Inches(5.15), Inches(1.0), WHITE,
         line=BLUE)
add_text(s, Inches(0.8), Inches(5.72), Inches(5.0), Inches(0.30),
         "지원 장비 · 포맷", size=11, bold=True, color=NAVY, font=FONT_TITLE)
add_text(s, Inches(0.8), Inches(6.02), Inches(5.0), Inches(0.65),
         "PNE (.cyc + SaveEnd.csv + channel_info.json + .log) · "
         "Toyo (TC CSV)\n신뢰성 .xls (Fasoo DRM) · .sch TC Plan · "
         "코인/3전극 (DB 적재 예정)",
         size=9, color=GREY_TXT, font=FONT_BODY)

add_text(s, Inches(6.3), Inches(1.3), Inches(6.7), Inches(0.5),
         "BDT 탭 구조 & 4대 미션 매핑",
         size=17, bold=True, color=NAVY, font=FONT_TITLE)

tabs = [
    ("현황/필터링",     "공통 기반",      "v4 (4/20) 완료",       GREEN_TAG),
    ("사이클데이터",    "공통 기반",      "파이프라인 리팩터 완",  GREEN_TAG),
    ("파형수정",        "미션 1 전처리",  "운영 중",               TEAL),
    ("셋 결과",         "수명·성능 요약", "운영 중",               TEAL),
    ("dVdQ 분석",       "미션 2",         "스무딩 1차, 2Q 매칭",   ORANGE_TAG),
    ("Eu 수명예측",     "미션 4 보조",    "운영 중",               TEAL),
    ("승인 수명예측",   "미션 4 (승인)",  "온도별 파라미터",       ORANGE_TAG),
    ("실사용 수명예측", "미션 4 (실사용)", "상위 유저 기준",         ORANGE_TAG),
    ("전기화학 Sim",    "미션 3",         "P2D/SPM/SPMe 운영",     TEAL),
]
ty = Inches(1.9)
for name, mission, status, color in tabs:
    add_rect(s, Inches(6.3), ty, Inches(6.7), Inches(0.48), WHITE,
             line=GREY_MUTED)
    add_text(s, Inches(6.45), ty, Inches(2.3), Inches(0.48),
             name, size=11, bold=True, color=NAVY, font=FONT_TITLE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(8.7), ty, Inches(2.0), Inches(0.48),
             mission, size=10, color=GREY_TXT, font=FONT_BODY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_status_pill(s, Inches(10.75), ty + Inches(0.08),
                    Inches(2.15), Inches(0.32), status, color)
    ty += Inches(0.53)

add_footer(s, 4, TOTAL)


# ============================ 미션 슬라이드 (v2: Science Note 3단) ============================
def mission_slide(num_str, main_title, subtitle, goal, bdt_role,
                  science_title, science_body,
                  achievements, plans, risk, page_no):
    s = prs.slides.add_slide(BLANK)
    add_confidential_stamp(s)
    add_title(s, num_str, main_title)

    add_text(s, Inches(0.45), Inches(0.92), Inches(10), Inches(0.30),
             subtitle, size=12, italic=True, color=BLUE, font=FONT_TITLE)

    # 좌측 : 목표 / BDT 역할 / Science Note (3단)
    add_rect(s, Inches(0.45), Inches(1.35), Inches(4.5), Inches(5.5), NAVY)

    # 1) 목표 (y 1.45~3.08)
    add_text(s, Inches(0.6), Inches(1.43), Inches(4.2), Inches(0.28),
             "목표", size=12, bold=True, color=LIGHT_BLUE, font=FONT_TITLE)
    add_rect(s, Inches(0.6), Inches(1.72), Inches(0.3), Inches(0.03), RED)
    add_text(s, Inches(0.6), Inches(1.78), Inches(4.2), Inches(1.28),
             goal, size=9.5, color=WHITE, font=FONT_BODY)

    # 2) BDT 역할 (y 3.15~4.52)
    add_text(s, Inches(0.6), Inches(3.13), Inches(4.2), Inches(0.28),
             "BDT 역할", size=12, bold=True, color=LIGHT_BLUE, font=FONT_TITLE)
    add_rect(s, Inches(0.6), Inches(3.42), Inches(0.3), Inches(0.03), RED)
    add_text(s, Inches(0.6), Inches(3.48), Inches(4.2), Inches(1.05),
             bdt_role, size=9.5, color=WHITE, font=FONT_BODY)

    # 3) Science Note (도메인 지식, y 4.58~6.85)
    add_rect(s, Inches(0.53), Inches(4.58), Inches(4.35), Inches(2.27), BLACK)
    add_text(s, Inches(0.63), Inches(4.65), Inches(4.15), Inches(0.28),
             f"◆ Science Note — {science_title}",
             size=11, bold=True, color=SCI_YELLOW, font=FONT_TITLE)
    add_rect(s, Inches(0.63), Inches(4.94), Inches(0.3), Inches(0.03),
             SCI_YELLOW)
    add_text(s, Inches(0.63), Inches(5.00), Inches(4.20), Inches(1.82),
             science_body, size=8.5, color=WHITE, font=FONT_MONO)

    # 우측 : 1Q 실적 + 2Q 계획
    add_rect(s, Inches(5.2), Inches(1.35), Inches(7.8), Inches(2.55),
             LIGHT_BLUE)
    add_text(s, Inches(5.4), Inches(1.45), Inches(7.5), Inches(0.4),
             "1분기 실적", size=15, bold=True, color=NAVY, font=FONT_TITLE)
    yy = Inches(1.9)
    for item in achievements:
        add_text(s, Inches(5.4), yy, Inches(0.3), Inches(0.3),
                 "✓", size=14, bold=True, color=GREEN_TAG, font=FONT_TITLE)
        add_text(s, Inches(5.7), yy, Inches(7.2), Inches(0.3), item,
                 size=10.5, color=GREY_TXT, font=FONT_BODY)
        yy += Inches(0.36)

    add_rect(s, Inches(5.2), Inches(4.0), Inches(7.8), Inches(2.55), SAND)
    add_text(s, Inches(5.4), Inches(4.1), Inches(7.5), Inches(0.4),
             "2분기 계획", size=15, bold=True, color=NAVY, font=FONT_TITLE)
    yy = Inches(4.55)
    for item in plans:
        add_text(s, Inches(5.4), yy, Inches(0.3), Inches(0.3),
                 "▶", size=11, bold=True, color=ORANGE_TAG, font=FONT_TITLE)
        add_text(s, Inches(5.7), yy, Inches(7.2), Inches(0.3), item,
                 size=10.5, color=GREY_TXT, font=FONT_BODY)
        yy += Inches(0.36)

    if risk:
        add_rect(s, Inches(5.2), Inches(6.65), Inches(7.8), Inches(0.4),
                 WHITE, line=RED)
        add_text(s, Inches(5.35), Inches(6.65), Inches(1.0), Inches(0.4),
                 "⚠ Risk", size=10, bold=True, color=RED, font=FONT_TITLE,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(6.3), Inches(6.65), Inches(6.65), Inches(0.4),
                 risk, size=9.5, color=GREY_TXT, font=FONT_BODY,
                 anchor=MSO_ANCHOR.MIDDLE)

    add_footer(s, page_no, TOTAL)
    return s


# ============================ S5. 미션 1 — 소재 DB ============================
mission_slide(
    num_str="3",
    main_title="미션 ① 소재 전기화학 물성 DB화",
    subtitle="코인 · 3전극 측정 → PostgreSQL 중앙 DB → 다중 모델 SELECT 비교",
    goal=("· 수명 raw 폴더 분산 → 매번 재파싱\n"
          "· Toyo 10 폴더 비교 : 20~50 초\n"
          "· 신뢰성 .xls : Fasoo DRM (xlwings COM)\n"
          "→ DB 적재 후 SELECT 1회로 다중 모델 비교"),
    bdt_role=("· PostgreSQL 스키마 설계\n"
              "  (product / test_group / cycle_summary)\n"
              "· db_connector 모듈 (SELECT-only)\n"
              "· DB 서버 : NAS / 별도 → 다수 동접"),
    science_title="GITT · 율별 충방전 · OCP",
    science_body=(
        "◎ GITT 확산계수 (고체상 Fick):\n"
        "  D_s = (4/πτ)·(mV_M/MS)²·(ΔE_s/ΔE_t)²\n"
        "  → SPM / SPMe 핵심 파라미터\n"
        "◎ 율별 충방전 → 분극 η · i₀ 추출:\n"
        "  η = (RT/αF)·ln(i/i₀)  [Tafel]\n"
        "  → Butler-Volmer 교환전류밀도\n"
        "◎ Half coin (vs Li):\n"
        "  단극 OCP U(x) · 이론 capacity\n"
        "◎ 3-electrode (Li ref.):\n"
        "  풀셀 중 양·음극 전위 개별 측정"),
    achievements=[
        "물성 파라미터 / DB 스키마 1차 설계 (3/20)",
        "코인 · 3전극 측정 SOP 문서화 (3/22)",
        "신뢰성 .xls → CSV 변환 파이프라인 (DRM 대응)",
    ],
    plans=[
        "BDT 내 DB 탭 신설 (읽기 전용 클라이언트)",
        "업체 수급 물성값 → DB 1차 적재",
        "Q8 main/sub 스펙시트 입수 후 반영",
    ],
    risk="DRM 환경 .xlsx 차단 → 번들 txt + 차트 PNG 재설계 회피 (4/15~20 개편 완료)",
    page_no=5,
)


# ============================ S6. 미션 2 — dV/dQ ============================
mission_slide(
    num_str="4",
    main_title="미션 ② dV/dQ 양음극 분리 (열화 모드)",
    subtitle="풀셀 dV/dQ 를 양·음극 성분으로 분리 → LLI (Li 재고손실) / LAM (활물질손실) 정량화",
    goal=("· 셀 용량 bar 모델 ↔ slip (Li 재고)\n"
          "· 열화 모드 규명\n"
          "   LLI : SEI 성장 · Li plating\n"
          "   LAM : 활물질 crack · delamination\n"
          "→ 셀 설계 / 수명 예측 피드백"),
    bdt_role=("· dVdQ 분석 탭\n"
              "· 풀셀 곡선 스무딩 + 단극 OCV 매칭\n"
              "  (Gr / Si / NMC)\n"
              "· 열화 셀 slip · mass → LAM · LLI 분리"),
    science_title="dV/dQ 피크 해석 (열화 시그니처)",
    science_body=(
        "◎ dV/dQ = OCV 기울기 (V/Ah)\n"
        "  피크 = 상전이 (stage transition)\n"
        "◎ Graphite : LiC₆ ↔ LiC₁₂ ↔ LiC₂₇\n"
        "  → SOC별 특징 피크 3~4 개\n"
        "◎ 열화 시그니처:\n"
        "  • 피크 이동 (slip)  → LLI\n"
        "  • 피크 높이 감소    → LAM_NE\n"
        "  • 양극 피크 수축    → LAM_PE\n"
        "  • 피크 폭 변화      → R_ct 증가\n"
        "◎ 3-bar 모델:\n"
        "  Q_cell ↔ Q_pos·m_pos ↔ Q_neg·m_neg"),
    achievements=[
        "dVdQ 스무딩 로직 이식 (규섭님 주도)",
        "단위 음극 (Gr / Si) OCV 모델 구축",
        "M2 양음극 코인셀 data @ S.LSI 1차 확보",
    ],
    plans=[
        "풀셀 ↔ 단극 피크 슬라이딩 · 매칭 알고리즘",
        "핵심 열화 모드 (LLI · LAM_PE · LAM_NE) 자동 추출",
        "3Q : Tool 내부 dVdQ 분석 탭 최종 통합",
        "4Q : 다수 사이클 열화 트렌드 자동 분석",
    ],
    risk="양음극 OCP profile 별도 확보 필수 — M2 완료, M1 / Q8 순차 진행",
    page_no=6,
)


# ============================ S7. 미션 3 — 성능 Sim ============================
mission_slide(
    num_str="5",
    main_title="미션 ③ 성능 시뮬레이션 (전기화학)",
    subtitle="PyBaMM 기반 P2D · SPMe · SPM — 온도 · 충전 조건별 성능 예측 + 리튬 석출 risk",
    goal=("· 설계 경쟁력 강화 · 개발 지원\n"
          "· 모델 선택 (복잡도 ↔ 정확도)\n"
          "   P2D / DFN  (고정밀)\n"
          "   SPMe · SPM (경량 · ROM)\n"
          "· Embedded ROM 타겟 : BMS 탑재"),
    bdt_role=("· 전기화학 시뮬레이션 탭 (Chen2020)\n"
              "· 운영 : CC/CCCV · GITT · HPPC · 커스텀\n"
              "· 출력 : 양/음극 OCP · 표면 Li · η_neg\n"
              "· 4/19 OSS 14종 비교 → PyBOP 통합"),
    science_title="DFN · Butler-Volmer · Li plating",
    science_body=(
        "◎ DFN (Doyle-Fuller-Newman, P2D):\n"
        "  고체  ∂c/∂t = D_s·∇²c   [Fick]\n"
        "  전해액  Nernst-Planck + 전하보존\n"
        "  계면  i = i₀[e^(αFη/RT) − e^(-(1-α)Fη/RT)]\n"
        "        (Butler-Volmer)\n"
        "◎ Li plating 조건:\n"
        "  η_neg = φ_s − φ_l − U_neg < 0 V\n"
        "  → Safety margin ≥ 0.02 V (고율)\n"
        "◎ Arrhenius 온도 의존성:\n"
        "  k(T) = A·exp(−E_a/RT)\n"
        "  SEI E_a ≈ 20~80 kJ/mol"),
    achievements=[
        "오픈소스 기반 모델 UI 구축 (탭 단위 통합)",
        "충 · 방전 / GITT / HPPC 시뮬레이션 (문헌 기반)",
        "전기화학 OSS 14종 비교 딥서치 (4/19) — PyBOP 통합 방향",
        "PyBaMM 변수 전수 조사 PPT 확보 (4/11)",
    ],
    plans=[
        "기본 모델 캘리브레이션 (실측 vs Sim 편차 ≤ 3 %)",
        "Arrhenius 온도 의존성 (저/고온) 부여",
        "PyBOP (BSD-3) 통합 — 셀별 파라미터 역추정",
        "3Q : 급속 스텝 → Li plating risk 도출 SW",
    ],
    risk="OSS 라이선스 경계 — cideMOD/CIDETEC PINN = AGPL (직접 포함 X), Wen PINN LICENSE 미명시",
    page_no=7,
)


# ============================ S8. 미션 4 — 수명 예측 ============================
mission_slide(
    num_str="6",
    main_title="미션 ④ 수명 예측 (Empirical 모델)",
    subtitle="평가환경 matrix (온도 · C-rate · 하한 SOC · 저장) → 수명 파라미터 → EOL 자동 예측",
    goal=("· 상품화(승인) 수명 — PF / 업체별 trend\n"
          "· 실사용 수명 — 상위 유저(고온·만충)\n"
          "· 평가환경 5축 :\n"
          "   온도 / 복합(충전 후 5·12hr)\n"
          "   방전 C-rate / 하한 SOC / 저장"),
    bdt_role=("· 승인 · 실사용 수명예측 탭\n"
              "· folderble → 파라미터 자동 도출\n"
              "· 1,500일 capacity · swelling 예측\n"
              "· EOL (SOH 80 %) 자동 계산"),
    science_title="Capacity Fade · Calendar/Cycle 분리",
    science_body=(
        "◎ 통합 fade 모델 (BDT):\n"
        "  Q(n) = 1 − a·n^b − c·exp(d·(n−e))\n"
        "       └ power law ┘└ knee exp ┘\n"
        "  · b ≈ 0.5 → SEI diffusion-lim.\n"
        "  · exp 항  → plating · crack 가속\n"
        "◎ Calendar fade (저장):\n"
        "  ΔQ_cal = k₀·√t · exp(−E_a/RT)·f(SOC)\n"
        "◎ Cycle fade (운행):\n"
        "  ΔQ_cyc = g(DOD, C_chg, V_min)·n\n"
        "◎ Si 특수성 (2Q 계수 추가):\n"
        "  부피팽창 ~300 % → SEI 재생 반복\n"
        "  저전위(V_min) stress 민감 큼\n"
        "  → 하한전압 열화계수 필수"),
    achievements=[
        "온도별 열화 파라미터 도출 (folderble)",
        "Cell / SET Capacity · Swelling 1,500일 예측 곡선 구현",
    ],
    plans=[
        "Si 하한 전압 열화 계수 추가 (Si 고함량 편차 축소)",
        "승인 / 실사용 시나리오 EOL 예측 기능화",
        "온도별 수명 데이터 수급 (23 · 35 · 45 ℃)",
        "3Q : 다조건 비교 + 수명 예측 자동화 리포트",
    ],
    risk="온도별 수명 raw 수급 선결 — 미수급 시 2Q 캘리브 ↔ 4Q 트렌드 지연",
    page_no=8,
)


# ============================ S9. 4월 개발 실적 ============================
s = prs.slides.add_slide(BLANK)
add_confidential_stamp(s)
add_title(s, "7", "4월 개발 실적 하이라이트")
add_text(s, Inches(0.45), Inches(0.92), Inches(10), Inches(0.30),
         "162 커밋 / 4개 카테고리 — 공통 기반 · 분석 · DRM · 테스트",
         size=13, italic=True, color=BLUE, font=FONT_TITLE)

cat_defs = [
    ("공통 기반 강화", "6건", NAVY,
     ["사이클 파이프라인 6단계 리팩터 (4/12)",
      "Logical cycle 3단 (TotlCycle↔논리↔OriCyc)",
      "Unified Profile (PNE/Toyo 공통)",
      "현황 탭 분류 v4 (JSON 1순위 + 교차검증)",
      "TC Plan .sch 기반 — Tier 2 recall 98.89 %",
      "ECT path 재정비 (Rest 체크, TC 힌트)"]),
    ("분석 기능", "5건", TEAL,
     ["프로파일 옵션 재설계 + 프리셋 + DOD축",
      "히스테리시스 major 임계값 0.98 상향",
      "F1 GITT D_s 심화 분석 도구",
      "ampworks F1~F8 차용 분석",
      "사이클 분류 9종 체계 (Phase 0/1/2)"]),
    ("DRM 대응", "환경", RED,
     ["Fasoo DRM : .xlsx/.pptx/.png/.csv 차단",
      "통과 가능 : 첫 줄 공란 .txt 전용",
      "--pack / --unpack OOXML base64",
      "--to-xlsx / pdf / docx / pptx 양방향",
      "교훈 : 사외 COM 에러 = 환경 제약"]),
    ("테스트 커버리지", "3계층", GREEN_TAG,
     ["E2E 자동화 (사이클데이터 탭 수동루틴 재현)",
      "Plot 검증 Layer 4a (픽셀 단위 비교)",
      "TC Plan / Tier 2 / GUI 스모크 통합 검증"]),
]

positions = [
    (Inches(0.45), Inches(1.35)),
    (Inches(6.90), Inches(1.35)),
    (Inches(0.45), Inches(4.3)),
    (Inches(6.90), Inches(4.3)),
]
for (cat, stat, color, items), (x, y) in zip(cat_defs, positions):
    add_rect(s, x, y, Inches(6.0), Inches(2.8), WHITE, line=GREY_MUTED)
    add_rect(s, x, y, Inches(0.2), Inches(2.8), color)
    add_text(s, x + Inches(0.35), y + Inches(0.1), Inches(3.5), Inches(0.45),
             cat, size=16, bold=True, color=NAVY, font=FONT_TITLE)
    add_text(s, x + Inches(4.0), y + Inches(0.05), Inches(1.9), Inches(0.85),
             stat, size=36, bold=True, color=color,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP, font=FONT_TITLE)
    add_rect(s, x + Inches(0.35), y + Inches(0.65), Inches(5.4),
             Inches(0.02), color)
    yy = y + Inches(0.78)
    for it in items:
        add_text(s, x + Inches(0.35), yy, Inches(0.2), Inches(0.30),
                 "•", size=12, bold=True, color=color, font=FONT_TITLE)
        add_text(s, x + Inches(0.55), yy, Inches(5.4), Inches(0.30),
                 it, size=10, color=GREY_TXT, font=FONT_BODY)
        yy += Inches(0.32)

add_footer(s, 9, TOTAL)


# ============================ S10. 2Q 과제 + 그룹장 협조 ============================
s = prs.slides.add_slide(BLANK)
add_confidential_stamp(s)
add_title(s, "8", "2분기 핵심 과제 & 그룹장 협조 요청")

add_rect(s, Inches(0.45), Inches(1.2), Inches(7.3), Inches(5.6), LIGHT_BLUE)
add_text(s, Inches(0.65), Inches(1.3), Inches(7.0), Inches(0.45),
         "2분기 핵심 과제",
         size=17, bold=True, color=NAVY, font=FONT_TITLE)

add_text(s, Inches(0.65), Inches(1.85), Inches(7.0), Inches(0.35),
         "반드시 끝내야 하는 것",
         size=13, bold=True, color=RED, font=FONT_TITLE)
must_items = [
    ("①", "DB UI 연동",       "미션 1",  "BDT 내 DB 탭 · 읽기 전용 클라이언트"),
    ("②", "풀셀 ↔ 단극 매칭",  "미션 2",  "피크 슬라이딩 · LLI/LAM 자동 추출"),
    ("③", "베이스 모델 캘리브", "미션 3",  "실측 vs Sim 편차 ≤ 3 %"),
    ("④", "Si 하한 전압 계수", "미션 4",  "Si 고함량 셀 예측 오차 축소"),
]
yy = Inches(2.25)
for num, title, mis, detail in must_items:
    add_rect(s, Inches(0.65), yy, Inches(7.0), Inches(0.5), WHITE)
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.75), yy + Inches(0.09),
                           Inches(0.32), Inches(0.32))
    c.fill.solid(); c.fill.fore_color.rgb = RED
    c.line.fill.background(); c.shadow.inherit = False
    add_text(s, Inches(0.75), yy + Inches(0.09), Inches(0.32), Inches(0.32),
             num, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
    add_text(s, Inches(1.15), yy, Inches(2.3), Inches(0.5),
             title, size=12, bold=True, color=NAVY, font=FONT_TITLE,
             anchor=MSO_ANCHOR.MIDDLE)
    add_status_pill(s, Inches(3.45), yy + Inches(0.1),
                    Inches(0.85), Inches(0.3), mis, NAVY)
    add_text(s, Inches(4.40), yy, Inches(3.2), Inches(0.5),
             detail, size=10, color=GREY_TXT, font=FONT_BODY,
             anchor=MSO_ANCHOR.MIDDLE)
    yy += Inches(0.55)

add_text(s, Inches(0.65), yy + Inches(0.15), Inches(7.0), Inches(0.35),
         "병행 과제",
         size=13, bold=True, color=TEAL, font=FONT_TITLE)
yy += Inches(0.55)
parallel = [
    "Arrhenius 온도 의존성 (미션 3)",
    "승인 / 실사용 EOL 시나리오 (미션 4)",
    "온도별 수명 데이터 수급 (23/35/45℃)",
    "PyBOP 통합 기초 작업 (미션 3, 4 연계)",
]
for it in parallel:
    add_text(s, Inches(0.8), yy, Inches(0.2), Inches(0.3),
             "▸", size=11, bold=True, color=TEAL, font=FONT_TITLE)
    add_text(s, Inches(1.0), yy, Inches(6.5), Inches(0.3),
             it, size=10, color=GREY_TXT, font=FONT_BODY)
    yy += Inches(0.26)

add_rect(s, Inches(7.95), Inches(1.2), Inches(5.05), Inches(5.6), NAVY)
add_text(s, Inches(8.15), Inches(1.3), Inches(4.7), Inches(0.45),
         "그룹장 협조 요청",
         size=17, bold=True, color=LIGHT_BLUE, font=FONT_TITLE)
add_rect(s, Inches(8.15), Inches(1.75), Inches(0.5), Inches(0.04), RED)

requests = [
    ("1", "DB 서버 / NAS 할당",
     "미션 1 DB 파이프라인 선결 조건"),
    ("2", "온도별 수명 데이터 수급\n(23 · 35 · 45 ℃)",
     "미션 4 Si 계수 · EOL 검증"),
    ("3", "Q8 main / sub\n업체 spec sheet",
     "미션 1 DB 초기 적재"),
    ("4", "PF / SDI / ATL / COSMX\n셀 수급 일정",
     "2Q 실험 계획 확정 (dVdQ 포함)"),
]
yy = Inches(1.95)
for num, title, why in requests:
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.15), yy + Inches(0.1),
                           Inches(0.38), Inches(0.38))
    c.fill.solid(); c.fill.fore_color.rgb = RED
    c.line.fill.background(); c.shadow.inherit = False
    add_text(s, Inches(8.15), yy + Inches(0.1), Inches(0.38), Inches(0.38),
             num, size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=FONT_TITLE)
    add_text(s, Inches(8.65), yy + Inches(0.05), Inches(4.2), Inches(0.5),
             title, size=11, bold=True, color=WHITE, font=FONT_TITLE)
    add_text(s, Inches(8.65), yy + Inches(0.55), Inches(4.2), Inches(0.5),
             why, size=9, italic=True, color=LIGHT_BLUE, font=FONT_BODY)
    yy += Inches(1.18)

add_footer(s, 10, TOTAL)


# ─────────────────────────── 저장 ───────────────────────────
out_dir = Path(__file__).resolve().parent
out_path = out_dir / "260421_보고_그룹장_선행Lab_BDT_현황및계획.pptx"
prs.save(str(out_path))
print(f"saved: {out_path}")
print(f"slides: {len(prs.slides)}")
